from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
PINK = RGBColor(0xFF, 0x2E, 0x97)
LIME = RGBColor(0xB8, 0xFF, 0x00)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
YELLOW = RGBColor(0xFF, 0xE6, 0x00)
LAVENDER = RGBColor(0xC8, 0xA2, 0xFF)
CREAM = RGBColor(0xFF, 0xF8, 0xF0)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_PINK = RGBColor(0xCC, 0x00, 0x66)
DARK_CYAN = RGBColor(0x00, 0x99, 0xBB)
LIGHT_PINK = RGBColor(0xFF, 0xE0, 0xF0)
LIGHT_CYAN = RGBColor(0xE0, 0xF8, 0xFF)
LIGHT_LIME = RGBColor(0xF0, 0xFF, 0xE0)
LIGHT_LAVENDER = RGBColor(0xF0, 0xE8, 0xFF)
LIGHT_YELLOW = RGBColor(0xFF, 0xFB, 0xE0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK, icon="▸", line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(font_size * 0.6)
    return txBox


def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# ============================================================
# SLIDE 1 — TITLE / COVER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK)

# Decorative circles
add_circle(slide, Inches(-1), Inches(-1), Inches(3), RGBColor(0xFF, 0x2E, 0x97))
add_circle(slide, Inches(11.5), Inches(5.5), Inches(3), RGBColor(0x00, 0xD4, 0xFF))
add_circle(slide, Inches(10), Inches(-0.5), Inches(1.5), LAVENDER)
add_circle(slide, Inches(0.5), Inches(6), Inches(1.5), LIME)

# Semi-transparent overlay for text area
add_shape(slide, Inches(2), Inches(1.5), Inches(9.333), Inches(4.5), RGBColor(0x22, 0x22, 0x3A), PINK, Pt(3))

# Title
add_text_box(slide, Inches(2.5), Inches(1.8), Inches(8.333), Inches(1.2),
             "seekho", 72, PINK, True, PP_ALIGN.CENTER, "Segoe UI Black")

# Tagline
add_text_box(slide, Inches(2.5), Inches(3.2), Inches(8.333), Inches(0.8),
             "swap skills. grow together. level up ur life.", 28, WHITE, False, PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(2.5), Inches(4.2), Inches(8.333), Inches(0.8),
             "India's First Peer-to-Peer Skill Exchange Platform for Students", 20, CYAN, False, PP_ALIGN.CENTER)

# Bottom bar
add_text_box(slide, Inches(2.5), Inches(5.2), Inches(8.333), Inches(0.5),
             "Community Onboarding Deck  |  2025-26", 14, GRAY, False, PP_ALIGN.CENTER)

# Corner labels
add_text_box(slide, Inches(0.3), Inches(7), Inches(3), Inches(0.4),
             "seekho.app", 12, LAVENDER, False, PP_ALIGN.LEFT)
add_text_box(slide, Inches(10), Inches(7), Inches(3), Inches(0.4),
             "Built by Students, for Students", 12, LIME, False, PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 2 — THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PINK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "The Problem", 40, PINK, True, PP_ALIGN.LEFT, "Segoe UI Black")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
             "College students have skills to share but no platform to exchange them.", 20, GRAY, False)

# Problem cards
problems = [
    ("😤", "No Way to Discover Peers", "You never know the guy next to you\nin class knows Python, guitar, or editing", LIGHT_PINK, PINK),
    ("💸", "Expensive Courses", "Students pay thousands for skills their\npeers could teach for free", LIGHT_CYAN, CYAN),
    ("😴", "Boring Traditional Learning", "Classroom learning doesn't cover\npractical, trending skills", LIGHT_LIME, LIME),
    ("🤷", "No Community Platform", "WhatsApp groups are chaotic. No proper\nspace for skill exchange at colleges", LIGHT_LAVENDER, LAVENDER),
]

for i, (emoji, title, desc, bg, accent) in enumerate(problems):
    left = Inches(0.6 + i * 3.1)
    top = Inches(2.2)
    card = add_shape(slide, left, top, Inches(2.9), Inches(4.2), bg, accent, Pt(2.5))

    add_text_box(slide, left + Inches(0.3), top + Inches(0.3), Inches(2.3), Inches(0.8),
                 emoji, 48, DARK, False, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.2), Inches(2.3), Inches(0.7),
                 title, 18, DARK, True, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(2.1), Inches(2.5), Inches(1.8),
                 desc, 14, GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
             "Seekho solves ALL of this — for free.", 18, PINK, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3 — WHAT IS SEEKHO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "What is Seekho?", 40, CYAN, True, PP_ALIGN.LEFT, "Segoe UI Black")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "A platform where students teach what they know & learn what they don't.\nThink Tinder, but for skills. You match, you learn, you grow.", 20, DARK, False)

# Feature grid (2x3)
features = [
    ("🔍", "Smart Skill Matching", "AI matches you with peers\nwho have what you need", LIGHT_CYAN, CYAN),
    ("💬", "Real-Time Chat", "Message your matches instantly.\nSchedule sessions in-app.", LIGHT_PINK, PINK),
    ("📅", "Session Booking", "Book 1-on-1 or group sessions.\nTrack your learning journey.", LIGHT_LIME, LIME),
    ("⭐", "Aura Points", "Earn reputation points for\nteaching & being awesome.", LIGHT_LAVENDER, LAVENDER),
    ("🏆", "Rewards System", "Redeem points for perks,\nbadges, and recognition.", LIGHT_YELLOW, YELLOW),
    ("👥", "College Communities", "Join your college chapter.\nAttend events & meetups.", LIGHT_PINK, PINK),
]

for i, (emoji, title, desc, bg, accent) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.1)
    top = Inches(2.5 + row * 2.3)
    card = add_shape(slide, left, top, Inches(3.8), Inches(2.0), bg, accent, Pt(2))

    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.7), Inches(0.6),
                 emoji, 30, DARK, False, PP_ALIGN.LEFT)
    add_text_box(slide, left + Inches(0.9), top + Inches(0.2), Inches(2.7), Inches(0.5),
                 title, 18, DARK, True, PP_ALIGN.LEFT)
    add_text_box(slide, left + Inches(0.9), top + Inches(0.8), Inches(2.7), Inches(1.0),
                 desc, 14, GRAY, False, PP_ALIGN.LEFT)


# ============================================================
# SLIDE 4 — HOW IT WORKS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), LIME)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "How It Works", 40, RGBColor(0x55, 0x99, 0x00), True, PP_ALIGN.LEFT, "Segoe UI Black")

steps = [
    ("01", "Sign Up", "Create your account with\nyour college email.\nTakes 30 seconds.", PINK),
    ("02", "List Skills", "Add skills you can teach\n& skills you wanna learn.\nBe honest!", CYAN),
    ("03", "Get Matched", "Our algorithm finds the\nbest peers for you.\nSwipe & connect.", LIME),
    ("04", "Book a Sesh", "Schedule a session —\nonline or in-person.\nYour choice.", LAVENDER),
    ("05", "Learn & Earn", "Exchange skills, earn\naura points, get rewards.\nLevel up!", YELLOW),
]

for i, (num, title, desc, accent) in enumerate(steps):
    left = Inches(0.4 + i * 2.5)
    top = Inches(1.8)

    # Number circle
    circle = add_circle(slide, left + Inches(0.7), top, Inches(1.0), accent)
    add_text_box(slide, left + Inches(0.7), top + Inches(0.1), Inches(1.0), Inches(0.8),
                 num, 36, WHITE if accent != LIME and accent != YELLOW else DARK, True, PP_ALIGN.CENTER, "Segoe UI Black")

    # Arrow (except last)
    if i < 4:
        add_text_box(slide, left + Inches(1.8), top + Inches(0.2), Inches(0.7), Inches(0.6),
                     "→", 30, GRAY, False, PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(1.3), Inches(2.0), Inches(0.5),
                 title, 22, DARK, True, PP_ALIGN.CENTER)

    # Desc
    add_text_box(slide, left + Inches(0.1), top + Inches(1.9), Inches(2.2), Inches(2.0),
                 desc, 14, GRAY, False, PP_ALIGN.CENTER)

# Bottom CTA
add_shape(slide, Inches(3.5), Inches(5.8), Inches(6.333), Inches(1.2), PINK, None)
add_text_box(slide, Inches(3.5), Inches(5.9), Inches(6.333), Inches(0.5),
             "It's FREE. No catch. Just learn.", 24, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.5), Inches(6.4), Inches(6.333), Inches(0.4),
             "seekho.app", 16, LIGHT_PINK, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5 — COMMUNITY STRUCTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), LAVENDER)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Community Structure", 40, LAVENDER, True, PP_ALIGN.LEFT, "Segoe UI Black")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
             "Every college gets its own Seekho chapter. Here's how the hierarchy works:", 18, GRAY, False)

# Pyramid-style hierarchy
levels = [
    ("🏢  Seekho HQ (Core Team)", "Platform development, strategy, partnerships, overall management", Inches(3.5), Inches(6.333), PINK, WHITE),
    ("🌟  City Leads", "Manage all college chapters in a city (Pune, Mumbai, Nagpur, etc.)", Inches(2.8), Inches(7.733), CYAN, WHITE),
    ("🎯  Campus Ambassadors", "Lead the Seekho community at their specific college", Inches(2.1), Inches(9.133), LAVENDER, WHITE),
    ("👥  Community Members", "All students who sign up and participate in skill exchange", Inches(1.4), Inches(10.533), LIME, DARK),
]

for i, (title, desc, top, width, bg, text_color) in enumerate(levels):
    left = (SLIDE_W - width) / 2
    card = add_shape(slide, left, top, width, Inches(1.1), bg, None)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.1), width - Inches(0.6), Inches(0.45),
                 title, 18, text_color, True, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.55), width - Inches(0.6), Inches(0.45),
                 desc, 13, text_color if bg == LIME else RGBColor(0xFF, 0xDD, 0xEE), False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6 — BECOME A MEMBER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Become a Community Member", 40, CYAN, True, PP_ALIGN.LEFT, "Segoe UI Black")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "Joining Seekho is easy & free. Here's what you get:", 20, DARK, False)

# Left side — How to join
add_shape(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(5.0), LIGHT_CYAN, CYAN, Pt(2))
add_text_box(slide, Inches(1.0), Inches(2.2), Inches(5), Inches(0.5),
             "✅  How to Join", 24, CYAN, True)

join_steps = [
    "Visit seekho.app and sign up with your college email",
    "Complete your profile — add skills you know & want to learn",
    "Join your college's community chapter",
    "Start discovering peers and booking sessions",
    "Attend community events and meetups",
]
add_bullet_list(slide, Inches(1.0), Inches(2.9), Inches(5.0), Inches(3.8),
                join_steps, 15, DARK, "→")

# Right side — What you get
add_shape(slide, Inches(6.9), Inches(2.0), Inches(5.8), Inches(5.0), LIGHT_PINK, PINK, Pt(2))
add_text_box(slide, Inches(7.3), Inches(2.2), Inches(5), Inches(0.5),
             "🎁  What You Get", 24, PINK, True)

benefits = [
    "Free peer-to-peer skill exchange",
    "Access to your college's Seekho community",
    "Aura points & rewards for participation",
    "Certificate of participation for active members",
    "Network with skilled students across Maharashtra",
    "Exclusive invites to workshops & events",
]
add_bullet_list(slide, Inches(7.3), Inches(2.9), Inches(5.0), Inches(3.8),
                benefits, 15, DARK, "★")

# Bottom
add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
             "No fees. No hidden charges. Just sign up and start learning.", 16, GRAY, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7 — BECOME AN AMBASSADOR
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PINK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Become a Campus Ambassador", 40, PINK, True, PP_ALIGN.LEFT, "Segoe UI Black")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
             "Lead the Seekho movement at your college. Build community. Build your resume.", 20, DARK, False)

# 3 columns
cols = [
    ("Who Can Apply?", [
        "Currently enrolled students",
        "Active on social media",
        "Good communication skills",
        "Passionate about community",
        "Minimum 1 year left in college",
        "Leadership experience (a plus)",
    ], LIGHT_PINK, PINK),
    ("What You'll Do", [
        "Build & lead your college chapter",
        "Organize monthly events/meetups",
        "Recruit members & grow community",
        "Moderate & maintain quality",
        "Be the face of Seekho on campus",
        "Report to your City Lead",
    ], LIGHT_CYAN, CYAN),
    ("What You'll Get", [
        "Official Ambassador certificate",
        "LinkedIn recommendation letter",
        "Free Seekho Premium forever",
        "Event budget & Seekho merch",
        "Direct mentorship from founders",
        "Exclusive Ambassador network",
    ], LIGHT_LAVENDER, LAVENDER),
]

for i, (title, items, bg, accent) in enumerate(cols):
    left = Inches(0.5 + i * 4.2)
    card = add_shape(slide, left, Inches(2.0), Inches(3.9), Inches(5.0), bg, accent, Pt(2.5))
    add_text_box(slide, left + Inches(0.3), Inches(2.2), Inches(3.3), Inches(0.5),
                 title, 22, accent, True, PP_ALIGN.CENTER)
    add_bullet_list(slide, left + Inches(0.3), Inches(2.9), Inches(3.3), Inches(3.8),
                    items, 14, DARK, "▸")


# ============================================================
# SLIDE 8 — AMBASSADOR APPLICATION PROCESS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), YELLOW)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Ambassador Application Process", 40, RGBColor(0xCC, 0xA0, 0x00), True, PP_ALIGN.LEFT, "Segoe UI Black")

process = [
    ("📝", "Apply Online", "Fill the Ambassador application\nform on seekho.app/ambassador", LIGHT_PINK, PINK),
    ("📞", "Interview", "Quick 15-min video call with\nour community team", LIGHT_CYAN, CYAN),
    ("✅", "Selection", "Selected ambassadors get\nonboarding kit & training", LIGHT_LIME, LIME),
    ("🚀", "Launch", "Set up your college chapter\n& start building community", LIGHT_LAVENDER, LAVENDER),
]

for i, (emoji, title, desc, bg, accent) in enumerate(process):
    left = Inches(0.5 + i * 3.2)
    top = Inches(1.8)
    card = add_shape(slide, left, top, Inches(2.9), Inches(3.2), bg, accent, Pt(2))

    # Step number
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(2.5), Inches(0.6),
                 f"Step {i+1}", 14, accent, True, PP_ALIGN.LEFT)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.6), Inches(2.3), Inches(0.7),
                 emoji, 40, DARK, False, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.4), Inches(2.3), Inches(0.5),
                 title, 20, DARK, True, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(2.0), Inches(2.5), Inches(1.0),
                 desc, 14, GRAY, False, PP_ALIGN.CENTER)

# Timeline
add_shape(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.6), DARK, None)
add_text_box(slide, Inches(1.2), Inches(5.5), Inches(4), Inches(0.4),
             "⏱  TIMELINE", 16, PINK, True)
add_text_box(slide, Inches(1.2), Inches(5.9), Inches(10.5), Inches(0.4),
             "Application → Interview (within 3 days) → Selection (within 5 days) → Onboarding (1 week) → You're LIVE!", 15, WHITE, False)
add_text_box(slide, Inches(1.2), Inches(6.4), Inches(10.5), Inches(0.4),
             "Total time: ~2 weeks from application to running your chapter", 14, CYAN, False)


# ============================================================
# SLIDE 9 — EVENTS & ACTIVITIES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), LIME)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Events & Activities", 40, RGBColor(0x55, 0x99, 0x00), True, PP_ALIGN.LEFT, "Segoe UI Black")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
             "Seekho communities run regular events to keep the vibe alive:", 18, GRAY, False)

events = [
    ("🎤", "Skill Showcase", "Members demo their skills\nin front of the community.\nDiscover hidden talent!", LIGHT_PINK, PINK),
    ("🏫", "Workshops", "Hands-on workshops on\ntrending skills — coding,\ndesign, music, etc.", LIGHT_CYAN, CYAN),
    ("🤝", "Skill Exchange Meetups", "In-person peer matching.\nMeet, match, exchange.\nNetworking on steroids.", LIGHT_LIME, LIME),
    ("🏆", "Hack & Learn", "Mini hackathons where\nteams build projects using\nnewly learned skills.", LIGHT_LAVENDER, LAVENDER),
]

for i, (emoji, title, desc, bg, accent) in enumerate(events):
    left = Inches(0.5 + i * 3.2)
    top = Inches(2.0)
    card = add_shape(slide, left, top, Inches(2.9), Inches(3.5), bg, accent, Pt(2))

    add_text_box(slide, left + Inches(0.3), top + Inches(0.3), Inches(2.3), Inches(0.7),
                 emoji, 40, DARK, False, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.1), Inches(2.3), Inches(0.5),
                 title, 18, DARK, True, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(1.8), Inches(2.5), Inches(1.5),
                 desc, 14, GRAY, False, PP_ALIGN.CENTER)

# Inter-college
add_shape(slide, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.3), DARK, PINK, Pt(2))
add_text_box(slide, Inches(1.0), Inches(5.9), Inches(5), Inches(0.4),
             "🔥  Inter-College Events (Monthly)", 18, PINK, True)
add_text_box(slide, Inches(1.0), Inches(6.3), Inches(11), Inches(0.6),
             "Cross-campus skill battles  •  City-wide meetups  •  Online skill challenges  •  Leaderboard competitions  •  Annual Seekho Fest", 15, WHITE, False)


# ============================================================
# SLIDE 10 — AURA POINTS & REWARDS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), LAVENDER)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Aura Points & Rewards", 40, LAVENDER, True, PP_ALIGN.LEFT, "Segoe UI Black")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
             "Everything you do on Seekho earns you Aura — your reputation score.", 18, GRAY, False)

# How to earn
add_shape(slide, Inches(0.6), Inches(1.9), Inches(6.0), Inches(5.0), LIGHT_LAVENDER, LAVENDER, Pt(2))
add_text_box(slide, Inches(1.0), Inches(2.1), Inches(5.2), Inches(0.5),
             "✨  How to Earn Aura", 22, LAVENDER, True)

earn_items = [
    "Complete a skill session  →  +50 aura",
    "Get a 5-star review  →  +30 aura",
    "Teach a new skill  →  +40 aura",
    "Attend a community event  →  +20 aura",
    "Refer a friend who signs up  →  +25 aura",
    "Consistent weekly activity  →  +15 aura/week",
    "Win a skill challenge  →  +100 aura",
]
add_bullet_list(slide, Inches(1.0), Inches(2.8), Inches(5.2), Inches(3.8),
                earn_items, 14, DARK, "⚡")

# Rewards
add_shape(slide, Inches(6.9), Inches(1.9), Inches(6.0), Inches(5.0), LIGHT_YELLOW, YELLOW, Pt(2))
add_text_box(slide, Inches(7.3), Inches(2.1), Inches(5.2), Inches(0.5),
             "🎁  Redeem Rewards", 22, RGBColor(0xCC, 0xA0, 0x00), True)

reward_items = [
    "500 aura  →  Seekho stickers & badge",
    "1000 aura  →  Seekho merch (t-shirt, hoodie)",
    "2000 aura  →  Premium features unlocked",
    "3000 aura  →  Certificate of Excellence",
    "5000 aura  →  Seekho Spotlight feature",
    "10000 aura  →  Mentorship from industry pros",
    "Top aura/month  →  Featured on Seekho socials",
]
add_bullet_list(slide, Inches(7.3), Inches(2.8), Inches(5.2), Inches(3.8),
                reward_items, 14, DARK, "🎯")


# ============================================================
# SLIDE 11 — CODE OF CONDUCT (Brief)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PINK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Community Guidelines", 40, PINK, True, PP_ALIGN.LEFT, "Segoe UI Black")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "We keep it real, respectful, and safe. Here are the ground rules:", 18, GRAY, False)

# Do's
add_shape(slide, Inches(0.6), Inches(1.9), Inches(5.8), Inches(5.2), RGBColor(0xE8, 0xFF, 0xE8), RGBColor(0x00, 0xAA, 0x00), Pt(2))
add_text_box(slide, Inches(1.0), Inches(2.1), Inches(5), Inches(0.5),
             "✅  DO THIS", 24, RGBColor(0x00, 0xAA, 0x00), True)

dos = [
    "Be kind, respectful, and inclusive to everyone",
    "Show up on time for your sessions",
    "Give honest & constructive feedback",
    "Represent your skills accurately",
    "Report any issues or violations you see",
    "Help newcomers feel welcome",
    "Respect everyone's privacy",
]
add_bullet_list(slide, Inches(1.0), Inches(2.8), Inches(5.0), Inches(4.0),
                dos, 15, DARK, "👍")

# Don'ts
add_shape(slide, Inches(6.9), Inches(1.9), Inches(5.8), Inches(5.2), RGBColor(0xFF, 0xE8, 0xE8), RGBColor(0xCC, 0x00, 0x00), Pt(2))
add_text_box(slide, Inches(7.3), Inches(2.1), Inches(5), Inches(0.5),
             "🚫  DON'T DO THIS", 24, RGBColor(0xCC, 0x00, 0x00), True)

donts = [
    "Harass, bully, or discriminate anyone",
    "Use Seekho for dating or non-learning purposes",
    "Create fake accounts or manipulate reviews",
    "Spam, scam, or promote unauthorized stuff",
    "Share others' personal info without consent",
    "No-show without cancelling sessions",
    "Use the platform for academic dishonesty",
]
add_bullet_list(slide, Inches(7.3), Inches(2.8), Inches(5.0), Inches(4.0),
                donts, 15, DARK, "⛔")


# ============================================================
# SLIDE 12 — ROADMAP / VISION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)

add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "The Seekho Roadmap", 40, CYAN, True, PP_ALIGN.LEFT, "Segoe UI Black")

phases = [
    ("Phase 1", "NOW", "Launch in 5 pilot colleges\nacross Pune & Mumbai.\nBuild core community.", PINK),
    ("Phase 2", "3 MONTHS", "Expand to 25+ colleges\nacross Maharashtra.\nLaunch mobile app.", CYAN),
    ("Phase 3", "6 MONTHS", "Cover all major cities in\nMaharashtra. 10,000+\nactive students.", LIME),
    ("Phase 4", "1 YEAR", "Go pan-India.\nPartner with companies.\nLaunch premium features.", LAVENDER),
    ("Phase 5", "2 YEARS", "India's largest student\nskill network. 1 lakh+\nusers. Revenue positive.", YELLOW),
]

for i, (phase, time, desc, accent) in enumerate(phases):
    left = Inches(0.3 + i * 2.6)
    top = Inches(1.8)

    # Timeline bar
    if i < 4:
        add_rect(slide, left + Inches(1.3), top + Inches(0.35), Inches(1.8), Inches(0.1), accent)

    # Circle
    add_circle(slide, left + Inches(0.8), top, Inches(0.8), accent)
    add_text_box(slide, left + Inches(0.8), top + Inches(0.15), Inches(0.8), Inches(0.5),
                 str(i + 1), 24, WHITE if accent != LIME and accent != YELLOW else DARK, True, PP_ALIGN.CENTER)

    # Phase title
    add_text_box(slide, left + Inches(0.2), top + Inches(1.2), Inches(2.2), Inches(0.4),
                 phase, 18, accent, True, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(1.6), Inches(2.2), Inches(0.3),
                 time, 14, GRAY, False, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(2.1), Inches(2.4), Inches(2.0),
                 desc, 14, WHITE, False, PP_ALIGN.CENTER)

# Bottom vision statement
add_shape(slide, Inches(2), Inches(5.5), Inches(9.333), Inches(1.5), RGBColor(0x22, 0x22, 0x3A), PINK, Pt(2))
add_text_box(slide, Inches(2.5), Inches(5.7), Inches(8.333), Inches(0.5),
             "Our Vision", 22, PINK, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.5), Inches(6.2), Inches(8.333), Inches(0.6),
             "Every student in India has a skill to share. Seekho makes sure no skill goes unlearned.", 18, WHITE, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 13 — FAQ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "FAQs", 40, CYAN, True, PP_ALIGN.LEFT, "Segoe UI Black")

faqs = [
    ("Is Seekho free?", "Yes! Core features are 100% free. Premium features coming soon at student-friendly prices."),
    ("Who can join?", "Any college student in Maharashtra (expanding to all India soon). Recent grads within 2 years also welcome."),
    ("Is it safe?", "We have strict community guidelines, verified college emails, reporting systems, and active moderators."),
    ("Can I be an Ambassador?", "If you're a current student with 1+ year left, good communication skills, and love building community — YES!"),
    ("What skills can I exchange?", "Anything! Coding, design, music, photography, cooking, fitness, languages, public speaking — you name it."),
    ("Do I get a certificate?", "Yes! Active members get participation certificates. Ambassadors get official recognition letters."),
]

for i, (q, a) in enumerate(faqs):
    col = i % 2
    row = i // 2
    left = Inches(0.5 + col * 6.4)
    top = Inches(1.4 + row * 1.95)

    colors = [LIGHT_PINK, LIGHT_CYAN, LIGHT_LIME, LIGHT_LAVENDER, LIGHT_YELLOW, LIGHT_PINK]
    accents = [PINK, CYAN, LIME, LAVENDER, YELLOW, PINK]

    card = add_shape(slide, left, top, Inches(6.0), Inches(1.7), colors[i], accents[i], Pt(2))
    add_text_box(slide, left + Inches(0.3), top + Inches(0.15), Inches(5.4), Inches(0.4),
                 f"Q: {q}", 16, accents[i], True)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.65), Inches(5.4), Inches(0.9),
                 a, 13, DARK, False)


# ============================================================
# SLIDE 14 — CTA / CLOSING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)

# Decorative
add_circle(slide, Inches(-0.5), Inches(-0.5), Inches(2.5), PINK)
add_circle(slide, Inches(11.5), Inches(5.5), Inches(2.5), CYAN)
add_circle(slide, Inches(5.5), Inches(-0.5), Inches(1.5), LAVENDER)
add_circle(slide, Inches(0.5), Inches(6.5), Inches(1), LIME)

add_shape(slide, Inches(1.5), Inches(1), Inches(10.333), Inches(5.5), RGBColor(0x22, 0x22, 0x3A), PINK, Pt(3))

add_text_box(slide, Inches(2), Inches(1.3), Inches(9.333), Inches(1),
             "Ready to Join the Movement?", 44, WHITE, True, PP_ALIGN.CENTER, "Segoe UI Black")

add_text_box(slide, Inches(2), Inches(2.4), Inches(9.333), Inches(0.8),
             "Seekho is more than an app — it's a community of students\nwho believe in learning from each other.", 22, GRAY, False, PP_ALIGN.CENTER)

# CTA buttons
add_shape(slide, Inches(2.5), Inches(3.6), Inches(3.8), Inches(0.9), PINK, None)
add_text_box(slide, Inches(2.5), Inches(3.7), Inches(3.8), Inches(0.7),
             "🚀  Join as Member\nseekho.app/join", 18, WHITE, True, PP_ALIGN.CENTER)

add_shape(slide, Inches(7), Inches(3.6), Inches(3.8), Inches(0.9), CYAN, None)
add_text_box(slide, Inches(7), Inches(3.7), Inches(3.8), Inches(0.7),
             "🎯  Apply as Ambassador\nseekho.app/ambassador", 18, WHITE, True, PP_ALIGN.CENTER)

# Social
add_text_box(slide, Inches(2), Inches(4.9), Inches(9.333), Inches(0.5),
             "Follow us:  @seekho.app  on Instagram  |  LinkedIn  |  Twitter", 16, LAVENDER, False, PP_ALIGN.CENTER)

# Contact
add_text_box(slide, Inches(2), Inches(5.5), Inches(9.333), Inches(0.5),
             "📧  hello@seekho.app   |   📱  seekho.app", 16, GRAY, False, PP_ALIGN.CENTER)

# Bottom tagline
add_text_box(slide, Inches(2), Inches(6.8), Inches(9.333), Inches(0.5),
             "swap skills. grow together. ✨", 20, PINK, True, PP_ALIGN.CENTER)


# Save
output_path = r"e:\XAMPP\htdocs\skillswap\legal\Seekho_Community_Deck.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
